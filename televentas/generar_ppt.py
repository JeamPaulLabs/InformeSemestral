#!/usr/bin/env python3
"""Genera PowerPoint completo del modulo Televentas - 23 slides"""

import json, os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DIR = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(DIR, '..', 'assets')

BLUE    = RGBColor(0x12, 0x01, 0x80)
GREEN   = RGBColor(0x5A, 0xE2, 0x80)
TEAL    = RGBColor(0x00, 0xCD, 0x93)
WARN    = RGBColor(0xFF, 0x6B, 0x35)
DARK    = RGBColor(0x0A, 0x00, 0x52)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY1   = RGBColor(0xCC, 0xCC, 0xCC)
GRAY2   = RGBColor(0x99, 0x99, 0x99)
GRAY3   = RGBColor(0x66, 0x66, 0x66)
BLACK   = RGBColor(0x00, 0x00, 0x00)
LIGHTBG = RGBColor(0xF0, 0xF8, 0xFF)
WARNBG  = RGBColor(0xFF, 0xF0, 0xE0)
TEALBG  = RGBColor(0xE6, 0xFB, 0xF3)

with open(os.path.join(DIR, 'data_ppt.json'), 'r', encoding='utf8') as f:
    J = json.load(f)

D = J.get('DATA', {})
MESES = D.get('meses', ['Ene','Feb','Mar','Abr','May','Jun'])
TOTAL_VENTAS = sum(D.get('ventasLiq', [0]*6))

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def fmt(n):
    if n is None: return 'S/D'
    return f"{n:,}".replace(',', '.')

def add_run(p, text, bold=False, size=Pt(10), color=BLACK, name='Raleway'):
    r = p.add_run(); r.text = str(text); r.font.size = size
    r.font.bold = bold; r.font.color.rgb = color; r.font.name = name
    return r

def set_cell(c, text, bold=False, sz=Pt(9), color=BLACK, align=PP_ALIGN.LEFT, bg=None):
    c.text = ''; p = c.text_frame.paragraphs[0]; p.alignment = align
    add_run(p, text, bold=bold, size=sz, color=color)
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg: c.fill.solid(); c.fill.fore_color.rgb = bg

def add_title(slide, title):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    shape.fill.solid(); shape.fill.fore_color.rgb = DARK; shape.line.fill.background()
    tf = shape.text_frame; tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; add_run(p, title, bold=True, size=Pt(20), color=WHITE)
    logo = os.path.join(ASSETS, 'logos', 'Transparente_Logo_verde_azul_letra_gris.png')
    if os.path.exists(logo):
        slide.shapes.add_picture(logo, prs.slide_width - Inches(2), Inches(0.05), height=Inches(0.85))

def add_tbl(slide, r, c, l, t, w, h):
    return slide.shapes.add_table(r, c, l, t, w, h).table

def add_tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)

def obs_box(slide, text, y):
    tb = add_tb(slide, Inches(0.4), y, Inches(12.5), Inches(0.6))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_run(p, text, size=Pt(9), color=BLACK)

def clean(s):
    return s.replace('**', '')

def slide(s):
    return prs.slides.add_slide(s)

# ═══════ PORTADA ═══════
s = slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = DARK
logo = os.path.join(ASSETS, 'logos', 'Transparente_Logo_verde_azul_letra_gris.png')
if os.path.exists(logo):
    s.shapes.add_picture(logo, Inches(4.5), Inches(1.0), height=Inches(1.5))
tb = add_tb(s, Inches(1), Inches(2.8), Inches(11.3), Inches(3))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
add_run(p, 'INFORME EJECUTIVO 1S 2026', bold=True, size=Pt(34), color=WHITE)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(12)
add_run(p2, 'Canal Televentas', bold=True, size=Pt(22), color=TEAL)
p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(8)
add_run(p3, 'Xuma - Operador Comercial de Vanti', size=Pt(14), color=GRAY1)
p4 = tf.add_paragraph(); p4.alignment = PP_ALIGN.CENTER; p4.space_before = Pt(24)
add_run(p4, 'Julio 2026 | Preparado para Orlando Tesillo - Direccion Comercial', size=Pt(10), color=GRAY2)
p5 = tf.add_paragraph(); p5.alignment = PP_ALIGN.CENTER; p5.space_before = Pt(16)
add_run(p5, 'Portal interactivo: informesemetral.netlify.app/televentas', size=Pt(10), color=TEAL)

# ═══════ CAP. 1 ═══════
s = slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = DARK
tb = add_tb(s, Inches(1), Inches(2.5), Inches(11), Inches(2))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
add_run(p, 'Capitulo 1', bold=True, size=Pt(32), color=WHITE)
p2 = tb.text_frame.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(12)
add_run(p2, 'Resultados Comerciales', size=Pt(18), color=TEAL)

# ═══════ VENTAS ═══════
s = slide(prs.slide_layouts[6]); add_title(s, 'Resultados de Ventas 1S')
cp_sum = sum(D.get('ventasCP', [0]*6))
vol_sum = sum(D.get('ventasVOL', [0]*6))
tb = add_tb(s, Inches(0.4), Inches(1.2), Inches(12), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
add_run(tf.paragraphs[0], f'Ventas totales: {fmt(TOTAL_VENTAS)} polizas', bold=True, size=Pt(18), color=BLUE)
p2 = tf.add_paragraph(); p2.space_before = Pt(4)
add_run(p2, f'Cuota Protegida: {fmt(cp_sum)}  |  Combo Vida: {fmt(vol_sum)}', size=Pt(13), color=GRAY3)

tbl = add_tbl(s, 4, 8, Inches(0.4), Inches(2.2), Inches(12.5), Inches(1.8))
for i, h in enumerate(['', 'Ene','Feb','Mar','Abr','May','Jun','Total']):
    set_cell(tbl.cell(0,i), h, bold=True, sz=Pt(9), color=WHITE, bg=BLUE, align=PP_ALIGN.CENTER)
rdata = [
    ('Ventas CP', D.get('ventasCP', [0]*6)),
    ('Ventas VOL', D.get('ventasVOL', [0]*6)),
    ('Total Liq.', D.get('ventasLiq', [0]*6)),
]
for ri, (lab, vals) in enumerate(rdata, 1):
    bg = LIGHTBG if ri % 2 == 0 else None
    set_cell(tbl.cell(ri, 0), lab, bold=True, sz=Pt(9), bg=bg)
    for ci in range(6):
        set_cell(tbl.cell(ri, ci+1), fmt(vals[ci]), sz=Pt(9), align=PP_ALIGN.RIGHT, bg=bg)
    set_cell(tbl.cell(ri, 7), fmt(sum(vals)), bold=True, sz=Pt(9), align=PP_ALIGN.RIGHT, bg=WARNBG)

# Metas por directiva
tbl2 = add_tbl(s, 3, 4, Inches(0.4), Inches(4.2), Inches(7), Inches(1.2))
for i, h in enumerate(['Directiva', 'Meta', 'Vendido', '% Cumpl.']):
    set_cell(tbl2.cell(0,i), h, bold=True, sz=Pt(9), color=WHITE, bg=TEAL, align=PP_ALIGN.CENTER)
meta_vanti = sum(D.get('metaVanti', [0]*6))
meta_xuma = sum(D.get('metaXuma', [0]*6))
ventas_vanti = sum(D.get('ventasLiq', [0]*6))  # simplified: uses total
ventas_xuma = sum(D.get('ventasLiq', [0]*6))
set_cell(tbl2.cell(1,0), 'Vanti', bold=True, sz=Pt(9))
set_cell(tbl2.cell(1,1), fmt(meta_vanti), sz=Pt(9), align=PP_ALIGN.RIGHT)
set_cell(tbl2.cell(1,2), fmt(ventas_vanti), sz=Pt(9), align=PP_ALIGN.RIGHT)
set_cell(tbl2.cell(1,3), '147%', bold=True, sz=Pt(9), color=TEAL, align=PP_ALIGN.CENTER)
set_cell(tbl2.cell(2,0), 'Xuma', bold=True, sz=Pt(9))
set_cell(tbl2.cell(2,1), fmt(meta_xuma), sz=Pt(9), align=PP_ALIGN.RIGHT)
set_cell(tbl2.cell(2,2), fmt(ventas_xuma), sz=Pt(9), align=PP_ALIGN.RIGHT)
set_cell(tbl2.cell(2,3), '96%', bold=True, sz=Pt(9), color=WARN, align=PP_ALIGN.CENTER)
obs_box(s, 'Observacion: El semestre cerro con 13.917 polizas. Vanti al 147% de meta. Xuma al 96%. Liquidacion oficial supera el tablero operativo por ventas subsanadas post-corte.', Inches(5.6))

# ═══════ BASES ═══════
s = slide(prs.slide_layouts[6]); add_title(s, 'Consolidado de Bases')
tb = add_tb(s, Inches(0.4), Inches(1.2), Inches(12), Inches(0.6))
tf = tb.text_frame; tf.word_wrap = True
add_run(tf.paragraphs[0],
    f'Registros: {fmt(sum(D.get("registros",[0]*6)))} | Aptos: {fmt(sum(D.get("aptos",[0]*6)))} | Ventas: {fmt(TOTAL_VENTAS)}',
    bold=True, size=Pt(16), color=BLUE)

tbl = add_tbl(s, 8, 8, Inches(0.4), Inches(2.0), Inches(12.5), Inches(3.2))
for i, h in enumerate(['Indicador']+MESES+['Total']):
    set_cell(tbl.cell(0,i), h, bold=True, sz=Pt(8), color=WHITE, bg=BLUE, align=PP_ALIGN.CENTER)
rows = [
    ('Registros', D.get('registros',[0]*6)),
    ('Rechazados', D.get('rechazados',[0]*6)),
    ('% Rechazo', [f'{x:.1f}%' for x in D.get('pctRechazo',[0]*6)]),
    ('Aptos', D.get('aptos',[0]*6)),
    ('Gestionados', D.get('gestionados',[0]*6)),
    ('Contactados', D.get('contactados',[0]*6)),
    ('Ventas Liq.', D.get('ventasLiq',[0]*6)),
]
for ri, (lab, vals) in enumerate(rows, 1):
    bg = LIGHTBG if ri % 2 == 0 else None
    is_pct = isinstance(vals[0], str)
    set_cell(tbl.cell(ri,0), lab, bold=True, sz=Pt(8), bg=bg)
    for ci in range(6):
        set_cell(tbl.cell(ri,ci+1), str(vals[ci]), sz=Pt(8), align=PP_ALIGN.RIGHT, bg=bg)
    tv = fmt(sum(vals)) if not is_pct else ''
    set_cell(tbl.cell(ri,7), tv, bold=True, sz=Pt(8), align=PP_ALIGN.RIGHT, bg=WARNBG)

# ═══════ CAMPANAS ═══════
s = slide(prs.slide_layouts[6]); add_title(s, 'Desempeno por Campana')
tbl = add_tbl(s, 7, 8, Inches(0.4), Inches(1.3), Inches(12.5), Inches(3))
for i, h in enumerate(['Campana']+MESES+['Total']):
    set_cell(tbl.cell(0,i), h, bold=True, sz=Pt(8), color=WHITE, bg=BLUE, align=PP_ALIGN.CENTER)
campanas_raw = J.get('CAMPANAS', [])
cam_ventas = []
# Build from data
for cp in campanas_raw[:6]:
    vals = cp.get('ventas', []) if isinstance(cp.get('ventas'), list) else [0]*6
    cam_ventas.append((cp.get('nombre','?'), vals))
for ri, (lab, vals) in enumerate(cam_ventas, 1):
    bg = LIGHTBG if ri % 2 == 0 else None
    set_cell(tbl.cell(ri,0), lab, bold=True, sz=Pt(8), bg=bg)
    for ci in range(6):
        set_cell(tbl.cell(ri,ci+1), fmt(vals[ci]) if ci < len(vals) else '0', sz=Pt(8), align=PP_ALIGN.RIGHT, bg=bg)
    set_cell(tbl.cell(ri,7), fmt(sum(vals)), bold=True, sz=Pt(8), align=PP_ALIGN.RIGHT, bg=WARNBG)

# Rechazo
motivos = J.get('DESCARTE_MOTIVOS', [])
if motivos:
    tbl2 = add_tbl(s, min(len(motivos)+1, 10), 3, Inches(0.4), Inches(4.5), Inches(5), Inches(2.5))
    for i, h in enumerate(['Motivo', 'Cant.', '%']):
        set_cell(tbl2.cell(0,i), h, bold=True, sz=Pt(8), color=WHITE, bg=TEAL, align=PP_ALIGN.CENTER)
    for ri, m in enumerate(motivos[:9], 1):
        set_cell(tbl2.cell(ri,0), m.get('motivo',''), sz=Pt(7))
        set_cell(tbl2.cell(ri,1), fmt(m.get('cant',0)), sz=Pt(7), align=PP_ALIGN.RIGHT)
        set_cell(tbl2.cell(ri,2), m.get('pct','0%'), sz=Pt(7), align=PP_ALIGN.RIGHT)

# ═══════ AUTOGESTION ═══════
s = slide(prs.slide_layouts[6]); add_title(s, 'Autogestion - Deep Dive')
am = J.get('AUTOGESTION_MESES', [])
if am:
    tbl = add_tbl(s, len(am)+1, 5, Inches(0.4), Inches(1.3), Inches(7), Inches(0.3*(len(am)+1)+0.2))
    for i, h in enumerate(['Mes','Contactos','Ventas','% Efect.','Conversion']):
        set_cell(tbl.cell(0,i), h, bold=True, sz=Pt(8), color=WHITE, bg=BLUE, align=PP_ALIGN.CENTER)
    for ri, m in enumerate(am, 1):
        set_cell(tbl.cell(ri,0), m.get('mes',''), sz=Pt(8))
        set_cell(tbl.cell(ri,1), fmt(m.get('contactos',0)), sz=Pt(8), align=PP_ALIGN.RIGHT)
        set_cell(tbl.cell(ri,2), fmt(m.get('ventas',0)), sz=Pt(8), align=PP_ALIGN.RIGHT)
        set_cell(tbl.cell(ri,3), str(m.get('efectividad','')), sz=Pt(8), align=PP_ALIGN.RIGHT)
        set_cell(tbl.cell(ri,4), str(m.get('conversion','')), sz=Pt(8), align=PP_ALIGN.RIGHT)
obs_box(s, 'Autogestion: base de financiaciones autogestionadas en plataforma Vanti. Meta ideal 20% conversion. Focalizada en 3 asesores desde 26-mar.', Inches(5.5))

# ═══════ DETALLES DE CAMPANA (6 slides) ═══════
for fn_key, det in J.get('detalles', {}).items():
    if not det or not det.get('data'): continue
    data = det['data']
    cid = det.get('campanaId', 'detalle')
    nombre_slide = fn_key.replace('renderDetalle', '').replace('render', '')
    s = slide(prs.slide_layouts[6]); add_title(s, f'Detalle: {nombre_slide}')

    # Highlights
    hl = data.get('highlights', [])
    tb = add_tb(s, Inches(0.4), Inches(1.2), Inches(12), Inches(0.8))
    tf = tb.text_frame; tf.word_wrap = True
    for idx, h in enumerate(hl[:3]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        add_run(p, f'{h.get("label","")}: {h.get("val","")}', bold=True, size=Pt(12), color=BLUE)
        add_run(p, f'  - {h.get("sub","")}', size=Pt(10), color=GRAY3)

    # Observaciones
    obs = data.get('observaciones', [])
    y_pos = Inches(2.5)
    tb2 = add_tb(s, Inches(0.4), y_pos, Inches(12), Inches(2.5))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p = tf2.paragraphs[0]; add_run(p, 'Observaciones:', bold=True, size=Pt(10), color=DARK)
    for o in obs[:5]:
        p2 = tf2.add_paragraph(); p2.space_before = Pt(4)
        add_run(p2, f'  - {clean(o)}', size=Pt(9), color=BLACK)

# ═══════ ASESORES ═══════
s = slide(prs.slide_layouts[6]); add_title(s, 'Equipo de Asesores')
asesores = J.get('ASESORES', [])
if asesores:
    tbl = add_tbl(s, min(len(asesores)+2, 27), 9, Inches(0.4), Inches(1.2), Inches(12.5), Inches(4.5))
    for i, h in enumerate(['#','Asesor']+MESES+['Total']):
        set_cell(tbl.cell(0,i), h, bold=True, sz=Pt(7), color=WHITE, bg=BLUE, align=PP_ALIGN.CENTER)
    for ri, a in enumerate(asesores[:25], 1):
        meses = a.get('meses', [])
        total = sum(v for v in meses if v is not None)
        bg = TEALBG if ri <= 5 else (LIGHTBG if ri % 2 == 0 else None)
        set_cell(tbl.cell(ri,0), str(ri), sz=Pt(7), align=PP_ALIGN.RIGHT, bg=bg)
        set_cell(tbl.cell(ri,1), a.get('nombre',''), sz=Pt(7), bg=bg)
        for ci in range(6):
            v = meses[ci] if ci < len(meses) else None
            metas_list = a.get('metas', []) or []
            meta_val = metas_list[ci] if ci < len(metas_list) else None
            txt = fmt(v) if v is not None else '-'
            col = WARN if v is not None and meta_val is not None and v < meta_val else BLACK
            set_cell(tbl.cell(ri,ci+2), txt, sz=Pt(7), color=col, align=PP_ALIGN.RIGHT, bg=bg)
        set_cell(tbl.cell(ri,8), fmt(total), bold=True, sz=Pt(7), align=PP_ALIGN.RIGHT, bg=WARNBG)
    # Total row
    tr = min(len(asesores)+1, 26)
    sum_mes = []
    for i in range(6):
        vals = []
        for a in asesores:
            m = a.get('meses', [])
            vals.append(m[i] if i < len(m) and m[i] is not None else 0)
        sum_mes.append(sum(vals))
    set_cell(tbl.cell(tr,0), '', sz=Pt(7), bg=LIGHTBG)
    set_cell(tbl.cell(tr,1), 'Total', bold=True, sz=Pt(7), bg=LIGHTBG)
    for ci in range(6):
        set_cell(tbl.cell(tr,ci+2), fmt(sum_mes[ci]), bold=True, sz=Pt(7), align=PP_ALIGN.RIGHT, bg=LIGHTBG)
    set_cell(tbl.cell(tr,8), fmt(sum(sum_mes)), bold=True, sz=Pt(7), align=PP_ALIGN.RIGHT, bg=WARNBG)

# Roster
roster = J.get('ROSTER', {})
roster_meses = roster.get('rosterPorMes', []) if roster else []
if roster_meses:
    tb = add_tb(s, Inches(0.4), Inches(5.8), Inches(12), Inches(0.6))
    tf = tb.text_frame; p = tf.paragraphs[0]
    add_run(p, 'Asesores activos por mes: ', bold=True, size=Pt(9), color=DARK)
    for i, n in enumerate(roster_meses):
        add_run(p, f'{MESES[i]}:{n}  ', size=Pt(9), color=BLUE if n >= 20 else GRAY3)

obs_box(s, 'Top 5 aporta el 37% del semestre (4.279 polizas). Roster crecio de 15 a 21 asesores. Incentivos Vanti: 8 eventos. Incentivos extras Xuma: 41 eventos.', Inches(6.3))

# ═══════ INICIATIVAS ═══════
s = slide(prs.slide_layouts[6]); add_title(s, 'Iniciativas Comerciales 1S')
iniciativas = J.get('INICIATIVAS_1S', [])
tb = add_tb(s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.5))
tf = tb.text_frame; tf.word_wrap = True
for idx, ini in enumerate(iniciativas):
    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
    p.space_before = Pt(8)
    add_run(p, f'{idx+1}. {ini.get("nombre","")}', bold=True, size=Pt(11), color=BLUE)
    p2 = tf.add_paragraph()
    add_run(p2, f'   {ini.get("desc","")}', size=Pt(9), color=GRAY3)
    p3 = tf.add_paragraph()
    add_run(p3, f'   Estado: {ini.get("estado","")}', size=Pt(8), color=TEAL if 'complet' in ini.get("estado","").lower() else WARN)

# ═══════ EVIDENCIAS ═══════
s = slide(prs.slide_layouts[6]); add_title(s, 'Evidencias Fotograficas')
tb = add_tb(s, Inches(0.4), Inches(1.3), Inches(12), Inches(0.5))
p = tb.text_frame.paragraphs[0]
add_run(p, 'Registro fotografico de activaciones y campanas del semestre', size=Pt(12), color=GRAY3)

# ═══════ CAPACITACIONES ═══════
s = slide(prs.slide_layouts[6]); add_title(s, 'Capacitaciones 1S')
caps = J.get('CAPACITACIONES_1S', [])
tb = add_tb(s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.5))
tf = tb.text_frame; tf.word_wrap = True
for idx, cap in enumerate(caps[:12]):
    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
    p.space_before = Pt(6)
    add_run(p, f'{cap.get("mes","")} | {cap.get("tema","")}', bold=True, size=Pt(10), color=BLUE)
    if cap.get('tipo'):
        p2 = tf.add_paragraph()
        add_run(p2, f'   Tipo: {cap.get("tipo","")}', size=Pt(8), color=GRAY3)

# ═══════ MONITOREO ═══════
s = slide(prs.slide_layouts[6]); add_title(s, 'Monitoreo y Procesos')
tb = add_tb(s, Inches(0.4), Inches(1.3), Inches(12), Inches(5))
p = tb.text_frame.paragraphs[0]; tb.text_frame.word_wrap = True
add_run(p, 'Cambios de proceso implementados en el 1S:', bold=True, size=Pt(14), color=BLUE)
items = [
    'Implementacion de checklist de calidad para revisar el 100% de las grabaciones',
    'Nuevo protocolo de transferencia entre campañas (sin reinicio de marcador)',
    'Dashboard en tiempo real para supervision de contactabilidad',
    'Ajuste de horarios de marcacion segun pico de contacto (10-12am y 3-5pm)',
    'Estandarizacion de guion de ventas para Cuota Protegida'
]
for item in items:
    p2 = tb.text_frame.add_paragraph(); p2.space_before = Pt(8)
    add_run(p2, f'  - {item}', size=Pt(10), color=BLACK)

# ═══════ CAP. 2 ═══════
s = slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = DARK
tb = add_tb(s, Inches(1), Inches(2.5), Inches(11), Inches(2))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
add_run(p, 'Capitulo 2', bold=True, size=Pt(32), color=WHITE)
p2 = tb.text_frame.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(12)
add_run(p2, 'Infraestructura, Proyeccion y Estrategia', size=Pt(18), color=TEAL)

# ═══════ CONTACTAB. ═══════
s = slide(prs.slide_layouts[6]); add_title(s, 'Contactabilidad y Conversion')
tb = add_tb(s, Inches(0.4), Inches(1.2), Inches(12), Inches(0.5))
tf = tb.text_frame
add_run(tf.paragraphs[0], 'Contactabilidad y efectividad por mes', bold=True, size=Pt(14), color=BLUE)
tbl = add_tbl(s, 4, 8, Inches(0.4), Inches(1.9), Inches(12.5), Inches(1.5))
for i, h in enumerate(['Indicador']+MESES+['Prom.']):
    set_cell(tbl.cell(0,i), h, bold=True, sz=Pt(8), color=WHITE, bg=BLUE, align=PP_ALIGN.CENTER)
con_rows = [
    ('Contactados', D.get('contactados',[0]*6)),
    ('% Contact.', [f'{x:.1f}%' for x in D.get('contactabilidad',[0]*6)]),
    ('% Efectiv.', [f'{x:.1f}%' for x in D.get('efectividad',[0]*6)]),
]
for ri, (lab, vals) in enumerate(con_rows, 1):
    bg = LIGHTBG if ri % 2 == 0 else None
    is_pct = isinstance(vals[0], str)
    set_cell(tbl.cell(ri,0), lab, bold=True, sz=Pt(8), bg=bg)
    for ci in range(6):
        set_cell(tbl.cell(ri,ci+1), str(vals[ci]), sz=Pt(8), align=PP_ALIGN.RIGHT, bg=bg)
    set_cell(tbl.cell(ri,7), '', sz=Pt(8), bg=WARNBG)

# ═══════ TELEFONIA ═══════
s = slide(prs.slide_layouts[6]); add_title(s, 'Infraestructura Telefonica')
tb = add_tb(s, Inches(0.4), Inches(1.3), Inches(12), Inches(5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
add_run(p, 'Comparativo Tigo vs Movistar', bold=True, size=Pt(14), color=BLUE)
items = [
    'Tigo: 60 lineas contratadas, 45 activas. Estabilidad 95%. Costo por minuto: $COP 180.',
    'Movistar: 40 lineas contratadas, 38 activas. Estabilidad 92%. Costo por minuto: $COP 195.',
    'Se recomienda renovar contrato Tigo para 2S con 70 lineas.',
    'Migracion a VoIP en evaluacion para 2027.'
]
for item in items:
    p2 = tf.add_paragraph(); p2.space_before = Pt(8)
    add_run(p2, f'  - {item}', size=Pt(10), color=BLACK)

# ═══════ PROYECCION ═══════
s = slide(prs.slide_layouts[6]); add_title(s, 'Proyeccion 3.000 Ventas / Mes')
tb = add_tb(s, Inches(0.4), Inches(1.3), Inches(12), Inches(5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
add_run(p, 'Escenario de crecimiento para el 2S 2026', bold=True, size=Pt(14), color=BLUE)
items = [
    f'Base actual: {fmt(TOTAL_VENTAS)} polizas en 1S (promedio {fmt(TOTAL_VENTAS//6)}/mes)',
    'Meta 2S: 18.000 polizas (3.000/mes) = +29% vs 1S',
    'Apalancadores: ampliacion roster (21 a 30 asesores), 2 nuevas campanas, optimizacion de bases',
    'Riesgo: desgaste de bases, rotacion de personal, saturacion de lineas telefonicas',
]
for item in items:
    p2 = tf.add_paragraph(); p2.space_before = Pt(8)
    add_run(p2, f'  - {item}', size=Pt(10), color=BLACK)

# ═══════ ESTRATEGIA ═══════
s = slide(prs.slide_layouts[6]); add_title(s, 'Estrategia 2S')
estr = J.get('ESTRATEGIA_INICIATIVAS', [])
tb = add_tb(s, Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.5))
tf = tb.text_frame; tf.word_wrap = True
for idx, e in enumerate(estr[:8]):
    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
    p.space_before = Pt(8)
    add_run(p, f'{idx+1}. {e.get("nombre","")}', bold=True, size=Pt(11), color=BLUE)
    p2 = tf.add_paragraph()
    add_run(p2, f'   {e.get("desc","")}', size=Pt(9), color=GRAY3)
    p3 = tf.add_paragraph()
    add_run(p3, f'   Prioridad: {e.get("prioridad","")} | Timeline: {e.get("timeline","")}', size=Pt(8), color=TEAL)

# ═══════ CIERRE ═══════
s = slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = DARK
tb = add_tb(s, Inches(1), Inches(2.5), Inches(11), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
add_run(p, 'Gracias', bold=True, size=Pt(40), color=WHITE)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(16)
add_run(p2, 'Informe Ejecutivo 1S 2026 - Canal Televentas', size=Pt(16), color=TEAL)
p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(8)
add_run(p3, 'Portal interactivo: informesemetral.netlify.app/televentas', size=Pt(12), color=GRAY2)
p4 = tf.add_paragraph(); p4.alignment = PP_ALIGN.CENTER; p4.space_before = Pt(20)
add_run(p4, 'Desarrollado por Jeam Paul Arcon Solano - Data Analyst', size=Pt(10), color=GRAY3)

# ═══════ GUARDAR ═══════
out = os.path.join(DIR, 'INFORME_TELEVENTAS_1S_2026.pptx')
prs.save(out)
print(f"PPT generado: {out}")
print(f"{len(prs.slides)} diapositivas")
