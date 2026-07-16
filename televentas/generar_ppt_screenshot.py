#!/usr/bin/env python3
"""
generar_ppt_screenshot.py
Genera el PPT de Televentas capturando screenshots pixel-perfect de cada slide.
- Resolución: 1920×1080 escalado a 1280×720 canvas → screenshot de 1920px de ancho
  para máxima calidad (luego PPT los comprime al tamaño de slide).
- Sin capas de texto superpuestas: fidelidad 100% al diseño web.
- Solo imágenes full-slide + notas de relator con el título de cada slide.
"""

import os, time, threading
from http.server import HTTPServer, SimpleHTTPRequestHandler
from functools import partial
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from playwright.sync_api import sync_playwright

DIR = os.path.dirname(os.path.abspath(__file__))
DASH_DIR = os.path.join(DIR, '..')
OUT_DIR = os.path.join(DIR, 'ppt_slides')
os.makedirs(OUT_DIR, exist_ok=True)

# Proporciones idénticas a las diapositivas web (1280×720 → 16:9 = 13.333" × 7.5")
PPT_W = Inches(13.333)
PPT_H = Inches(7.5)

# Tab definitions (label, JS click)
TABS_VTAS = [
    ('Vanti', "document.querySelector('.vtas-tab[data-vtab=\"vanti\"]')?.click()"),
    ('Xuma',  "document.querySelector('.vtas-tab[data-vtab=\"xuma\"]')?.click()"),
]
TABS_ASESORES = [
    ('Vanti', "document.querySelector('.asesor-tab[data-astab=\"vanti\"]')?.click()"),
    ('Xuma',  "document.querySelector('.asesor-tab[data-astab=\"xuma\"]')?.click()"),
]
TABS_CONTACTAB = [
    ('Mes',     "document.querySelector('.contactab-tab[data-ctab=\"mes\"]')?.click()"),
    ('Campaña', "document.querySelector('.contactab-tab[data-ctab=\"campana\"]')?.click()"),
]
TABS_TELE = [
    ('Resumen', "document.querySelector('.tele-tab[data-telab=\"resumen\"]')?.click()"),
    ('Zonas',   "document.querySelector('.tele-tab[data-telab=\"zonas\"]')?.click()"),
]
TABS_PROY = [
    ('Cálculo',   "document.querySelector('.proy-tab[data-ptab=\"calculo\"]')?.click()"),
    ('Escenario', "document.querySelector('.proy-tab[data-ptab=\"escenario\"]')?.click()"),
]
TABS_STRAT = [
    ('Iniciativas', "document.querySelector('.strategy-tab[data-stab=\"iniciativas\"]')?.click()"),
    ('Cronograma',  "document.querySelector('.strategy-tab[data-stab=\"cronograma\"]')?.click()"),
    ('KPIs',        "document.querySelector('.strategy-tab[data-stab=\"kpis\"]')?.click()"),
]

# (slide_num, label, [(tab_label, tab_js)] or None)
SLIDES = [
    (0,  'Portada',              None),
    (1,  'Cap. 1 – Gestión del Canal', None),
    (2,  'Ventas',               TABS_VTAS),
    (3,  'Bases Operativas',     None),
    (4,  'Detalle Campañas',     None),
    (5,  'Autogestión',          None),
    (6,  'D. Bienvenida CP',     None),
    (7,  'D. Stock',             None),
    (8,  'D. Masiva',            None),
    (9,  'D. Clientes Satisfechos', None),
    (10, 'D. Microseguro',       None),
    (11, 'D. Cancelaciones',     None),
    (12, 'Asesores',             TABS_ASESORES),
    (13, 'Iniciativas 1S',       None),
    (14, 'Evidencias',           None),
    (15, 'Capacitaciones',       None),
    (16, 'Monitoreo de Calidad', None),
    (17, 'Cap. 2 – Indicadores Operativos', None),
    (18, 'Contactabilidad',      TABS_CONTACTAB),
    (19, 'Telefonía',            TABS_TELE),
    (20, 'Proyección 2S',        TABS_PROY),
    (21, 'Estrategia 2S',        TABS_STRAT),
    (22, 'Cierre',               None),
]

# ── HTTP Server ──────────────────────────────────────────────────────────────
class SilentHandler(SimpleHTTPRequestHandler):
    def log_message(self, *a): pass

server = HTTPServer(('localhost', 0), partial(SilentHandler, directory=DASH_DIR))
port = server.server_address[1]
print(f"Servidor HTTP: http://localhost:{port}")
t = threading.Thread(target=server.serve_forever, daemon=True)
t.start()
time.sleep(0.6)

# ── Build PPT ────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = PPT_W
prs.slide_height = PPT_H
blank_layout = prs.slide_layouts[6]

print("Capturando slides...")

with sync_playwright() as pw:
    browser = pw.chromium.launch(headless=True)
    # Usar viewport 1280×720 exactos para que el deck-engine escale 1:1
    context = browser.new_context(
        viewport={'width': 1280, 'height': 720},
        device_scale_factor=2,   # 2x para alta resolución en PPT
    )
    page = context.new_page()

    def capture_slide(num: int, label: str, tab_js: str | None = None):
        """Captura screenshot de la slide y lo guarda. Retorna la ruta."""
        safe_label = label.replace(' ', '_').replace('/', '-').replace('–', '-')
        out_path = os.path.join(OUT_DIR, f'sk_{num:02d}_{safe_label}.png')

        page.goto(
            f"http://localhost:{port}/televentas/index.html?slide={num}",
            wait_until='networkidle',
            timeout=20000,
        )
        # Esperar a que las animaciones terminen
        time.sleep(0.9)

        # Ocultar cromo de navegación y forzar escala 1:1
        page.evaluate("""() => {
            // Ocultar todos los botones flotantes de navegación
            ['next-btn', 'prev-btn', 'home-btn', 'nav', 'counter', 'progress'].forEach(id => {
                const el = document.getElementById(id);
                if (el) el.style.setProperty('display', 'none', 'important');
            });
            document.querySelectorAll('.progress-bar, #progress').forEach(el => {
                el.style.setProperty('display', 'none', 'important');
            });
            // Forzar que el scaler esté en posición 0,0 con escala 1
            const s = document.getElementById('scaler');
            if (s) {
                s.style.transform = 'scale(1)';
                s.style.transformOrigin = 'top left';
                s.style.left = '0';
                s.style.top = '0';
                s.style.position = 'absolute';
            }
            document.body.style.overflow = 'hidden';
            document.body.style.margin = '0';
            document.body.style.padding = '0';
        }""")
        time.sleep(0.2)

        # Clic en la pestaña si corresponde
        if tab_js:
            page.evaluate(tab_js)
            time.sleep(0.4)

        # Capturar únicamente el área 1280×720 del canvas
        page.screenshot(
            path=out_path,
            clip={'x': 0, 'y': 0, 'width': 1280, 'height': 720},
        )
        return out_path

    def add_ppt_slide(img_path: str, notes_text: str):
        """Agrega una diapositiva al PPT con la imagen a tamaño completo."""
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_path, 0, 0, PPT_W, PPT_H)
        # Notas de relator
        notes = slide.notes_slide
        tf = notes.notes_text_frame
        tf.text = notes_text

    for num, label, tabs in SLIDES:
        try:
            if tabs:
                for tab_label, tab_js in tabs:
                    full_label = f"{label} – {tab_label}"
                    img_path = capture_slide(num, full_label, tab_js)
                    add_ppt_slide(img_path, f"Slide {num} | {full_label}")
                    print(f"  ✓ {num:2d}  {full_label}")
            else:
                img_path = capture_slide(num, label)
                add_ppt_slide(img_path, f"Slide {num} | {label}")
                print(f"  ✓ {num:2d}  {label}")
        except Exception as e:
            print(f"  ✗ {num:2d}  {label}: ERROR → {e}")

    browser.close()

out_path = os.path.join(DIR, 'INFORME_TELEVENTAS_1S_2026.pptx')
prs.save(out_path)
print(f"\n✅ PowerPoint guardado en: {out_path}")
print(f"   Total diapositivas: {len(prs.slides)}")
