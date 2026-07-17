#!/usr/bin/env python3
"""Test: poner texto en posiciones FIJAS conocidas para ver si el scaling funciona"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

DIR = r'C:\01_Repositorios\13_Presentaciones\dash\televentas'
ss_path = os.path.join(DIR, 'ppt_slides', 'sk_02_Vanti.png')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s = prs.slides.add_slide(prs.slide_layouts[6])

if os.path.exists(ss_path):
    s.shapes.add_picture(ss_path, 0, 0, Inches(13.333), Inches(7.5))

# Text boxes at KNOWN canvas positions
test_positions = [
    (80, 104, "Pólizas liquidadas (1S completo)"),
    (80, 122, "13.917"),
    (470, 104, "Asesores en equipo"),
    (470, 122, "15 → 21 → 21"),
    (861, 104, "Brecha vs. meta 3.000/mes"),
    (861, 122, "+30 %"),
]

SX = 13.333 / 1280.0
SY = 7.5 / 720.0

for tx, ty, text in test_positions:
    left = tx * SX
    top = ty * SY
    
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(3), Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = RGBColor(18, 1, 128)

out = os.path.join(DIR, 'TEST_POSICIONES.pptx')
prs.save(out)
print(f"Saved: {out}")
print(f"CanvasPos: (80,104) -> PPT ({80*SX:.3f}\", {104*SY:.3f}\")")
print(f"CanvasPos: (470,104) -> PPT ({470*SX:.3f}\", {104*SY:.3f}\")")
