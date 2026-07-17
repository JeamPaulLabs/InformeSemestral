#!/usr/bin/env python3
"""PPT hibrido v2: screenshot con texto + capas editables exactas + tabs separados"""

import json, os, time, threading, re
from http.server import HTTPServer, SimpleHTTPRequestHandler
from functools import partial
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from playwright.sync_api import sync_playwright

DIR = os.path.dirname(os.path.abspath(__file__))
DASH_DIR = os.path.join(DIR, '..')
OUT_DIR = os.path.join(DIR, 'ppt_slides')
os.makedirs(OUT_DIR, exist_ok=True)

PPT_W = Inches(13.333)
PPT_H = Inches(7.5)
SX = 13.333 / 1280.0  # 0.010416
SY = 7.5 / 720.0

def px(x, y):
    return (x * SX, y * SY)

# ── Server ──
class H(SimpleHTTPRequestHandler):
    def log_message(self, *a): pass
server = HTTPServer(('localhost', 0), partial(H, directory=DASH_DIR))
port = server.server_address[1]
print(f"Server: http://localhost:{port}")
t = threading.Thread(target=server.serve_forever, daemon=True); t.start()
time.sleep(0.5)

prs = Presentation()
prs.slide_width = PPT_W
prs.slide_height = PPT_H

# ── Helper ──
def parse_color(s):
    m = re.match(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', str(s))
    return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3))) if m else None

# ── Tab definitions: (tab_click_js, label_suffix)
TABS = {
    'vtas': [('Vanti', "document.querySelector('.vtas-tab[data-vtab=\"vanti\"]')?.click()"),
             ('Xuma', "document.querySelector('.vtas-tab[data-vtab=\"xuma\"]')?.click()")],
    'asesores': [('Vanti', "document.querySelector('.asesor-tab[data-astab=\"vanti\"]')?.click()"),
                 ('Xuma', "document.querySelector('.asesor-tab[data-astab=\"xuma\"]')?.click()")],
    'contactab': [('Mes', "document.querySelector('.contactab-tab[data-ctab=\"mes\"]')?.click()"),
                  ('Campana', "document.querySelector('.contactab-tab[data-ctab=\"campana\"]')?.click()")],
    'telefonia': [('Resumen', "document.querySelector('.tele-tab[data-telab=\"resumen\"]')?.click()"),
                  ('Zonas', "document.querySelector('.tele-tab[data-telab=\"zonas\"]')?.click()")],
    'proyeccion': [('Calculo', "document.querySelector('.proy-tab[data-ptab=\"calculo\"]')?.click()"),
                   ('Escenario', "document.querySelector('.proy-tab[data-ptab=\"escenario\"]')?.click()")],
    'estrategia': [('Iniciativas', "document.querySelector('.strategy-tab[data-stab=\"iniciativas\"]')?.click()"),
                   ('Cronograma', "document.querySelector('.strategy-tab[data-stab=\"cronograma\"]')?.click()"),
                   ('KPIs', "document.querySelector('.strategy-tab[data-stab=\"kpis\"]')?.click()")],
}

# ── Slides list: (num, label, [(tab_suffix, tab_click_js), ...] or None)
SLIDES = [
    (0, 'Portada', None),
    (1, 'Cap. 1', None),
    (2, 'Ventas', TABS['vtas']),
    (3, 'Bases', None),
    (4, 'Campanas', None),
    (5, 'Autogestion', None),
    (6, 'D. Bienvenida', None),
    (7, 'D. Stock', None),
    (8, 'D. Masiva', None),
    (9, 'D. Satisfechos', None),
    (10, 'D. Microseguro', None),
    (11, 'D. Cancelaciones', None),
    (12, 'Asesores', TABS['asesores']),
    (13, 'Iniciativas', None),
    (14, 'Evidencias', None),
    (15, 'Capacitaciones', None),
    (16, 'Monitoreo', None),
    (17, 'Cap. 2', None),
    (18, 'Contactab.', TABS['contactab']),
    (19, 'Telefonia', TABS['telefonia']),
    (20, 'Proyeccion', TABS['proyeccion']),
    (21, 'Estrategia', TABS['estrategia']),
    (22, 'Cierre', None),
]

print("Capturando slides y extrayendo textos...")

with sync_playwright() as pw:
    browser = pw.chromium.launch(headless=True)
    context = browser.new_context(viewport={'width': 1280, 'height': 720})
    page = context.new_page()
    
    for num, label, tabs in SLIDES:
        try:
            page.goto(f"http://localhost:{port}/televentas/index.html?slide={num}", wait_until='networkidle', timeout=15000)
            time.sleep(0.8)
            
            # Force canvas to exact viewport position
            page.evaluate("""() => {
                const s = document.getElementById('scaler');
                if (s) { s.style.transform = 'scale(1)'; s.style.top = '0'; s.style.left = '0'; }
                document.body.style.margin = '0';
                document.body.style.overflow = 'hidden';
            }""")
            time.sleep(0.2)
            
            # Hide nav chrome
            page.evaluate("""() => {
                // Hide all floating navigation buttons
                ['next-btn', 'prev-btn', 'home-btn', 'nav'].forEach(id => {
                    const el = document.getElementById(id);
                    if (el) el.style.setProperty('display', 'none', 'important');
                });
                const p = document.querySelector('.progress-bar');
                if (p) p.style.setProperty('display', 'none', 'important');
            }""")
            time.sleep(0.1)
            
            # Determine tab variants
            variants = [('', None)]  # (suffix, click_js)
            if tabs:
                variants = [(f' {suffix}', js) for suffix, js in tabs]
            
            for v_suffix, v_click in variants:
                slide_label = label + v_suffix
                
                # Click tab if needed
                if v_click:
                    page.evaluate(v_click)
                    time.sleep(0.3)
                
                # Screenshot with text intact
                ss_path = os.path.join(OUT_DIR, f'slide_{num:02d}{v_suffix.replace(" ","_")}.png')
                page.screenshot(path=ss_path)
                
                # Extract text positions
                text_data = page.evaluate("""() => {
                    const results = [];
                    const active = document.querySelector('#scaler .slide.active') || document.querySelector('#scaler');
                    if (!active) return results;
                    
                    const allEls = active.querySelectorAll('*');
                    allEls.forEach(el => {
                        const style = window.getComputedStyle(el);
                        if (style.display === 'none' || style.visibility === 'hidden') return;
                        if (el.closest('svg')) return;
                        
                        let dt = '';
                        for (const c of el.childNodes) {
                            if (c.nodeType === 3) dt += c.textContent;
                        }
                        const text = dt.trim();
                        if (!text || text.length < 2) return;
                        
                        const r = el.getBoundingClientRect();
                        if (r.width < 3 || r.height < 3) return;
                        
                        results.push({
                            text,
                            x: r.left, y: r.top,
                            w: r.width, h: r.height,
                            fontSize: parseFloat(style.fontSize) || 10,
                            bold: style.fontWeight === 'bold' || parseInt(style.fontWeight) >= 600,
                            color: style.color,
                            textAlign: style.textAlign || 'left'
                        });
                    });
                    return results;
                }""")
                
                # Deduplicate overlapping regions
                text_data.sort(key=lambda t: -len(t['text']))
                filtered = []
                regions = []
                for t in text_data:
                    rx, ry, rw, rh = t['x'], t['y'], t['w'], t['h']
                    contained = False
                    for ux, uy, uw, uh in regions:
                        if rx >= ux and ry >= uy and rx+rw <= ux+uw and ry+rh <= uy+uh:
                            contained = True; break
                    if not contained:
                        filtered.append(t)
                        regions.append((rx, ry, rw, rh))
                
                # Build PPT slide
                s = prs.slides.add_slide(prs.slide_layouts[6])
                
                # Screenshot as background
                if os.path.exists(ss_path):
                    s.shapes.add_picture(ss_path, 0, 0, PPT_W, PPT_H)
                
                # Editable text boxes at exact positions
                for t in filtered:
                    tx, ty, tw, th = t['x'], t['y'], t['w'], t['h']
                    if tw < 5 or th < 3: continue
                    
                    left, top = px(tx, ty)
                    width = max(tw * SX, Inches(0.25))
                    height = max(th * SY, Inches(0.12))
                    
                    tb = s.shapes.add_textbox(left, top, width, height)
                    tf = tb.text_frame
                    tf.word_wrap = False
                    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
                    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
                    
                    p = tf.paragraphs[0]
                    p.alignment = {'left': PP_ALIGN.LEFT, 'right': PP_ALIGN.RIGHT,
                                   'center': PP_ALIGN.CENTER}.get(t['textAlign'], PP_ALIGN.LEFT)
                    
                    fs = Pt(min(t['fontSize'] * 0.75, 16))
                    r = p.add_run()
                    r.text = t['text']
                    r.font.size = fs
                    r.font.bold = t['bold']
                    r.font.name = 'Raleway'
                    
                    clr = parse_color(t['color'])
                    if clr: r.font.color.rgb = clr
                    
                    tb.fill.background()
                
                print(f"  {num:2d} {slide_label:25s} -> {len(filtered):3d} textos")
                
        except Exception as e:
            print(f"  {num:2d} {label}: ERROR {e}")
    
    browser.close()

out = os.path.join(DIR, 'INFORME_TELEVENTAS_HIBRIDO.pptx')
prs.save(out)
print(f"\nPPT: {out}")
print(f"{len(prs.slides)} slides")
