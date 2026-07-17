#!/usr/bin/env python3
"""Genera PPT hibrido: fondo screenshot + capas editables de datos clave"""

import json, os, sys, time, threading
from http.server import HTTPServer, SimpleHTTPRequestHandler
from functools import partial
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from playwright.sync_api import sync_playwright

DIR = os.path.dirname(os.path.abspath(__file__))
DASH_DIR = os.path.join(DIR, '..')
OUT_DIR = os.path.join(DIR, 'ppt_slides')
os.makedirs(OUT_DIR, exist_ok=True)

PPT_W = Inches(13.333)
PPT_H = Inches(7.5)

BLUE = RGBColor(0x12, 0x01, 0x80)
TEAL = RGBColor(0x00, 0xCD, 0x93)
DARK = RGBColor(0x0A, 0x00, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY1 = RGBColor(0xCC, 0xCC, 0xCC)
GRAY3 = RGBColor(0x66, 0x66, 0x66)
BLACK = RGBColor(0x00, 0x00, 0x00)
OVERLAY_BG = RGBColor(0xF5, 0xF0, 0xFF)  # subtle purple tint

with open(os.path.join(DIR, 'data_ppt.json'), 'r', encoding='utf8') as f:
    J = json.load(f)
D = J.get('DATA', {})
MESES = D.get('meses', ['Ene','Feb','Mar','Abr','May','Jun'])

def fmt(n):
    if n is None: return 'S/D'
    return f"{n:,}".replace(',', '.')

def clean(s):
    return s.replace('**', '')

SLIDES = [
    (0, 'Portada'),
    (1, 'Cap. 1'),
    (2, 'Ventas'),
    (3, 'Bases'),
    (4, 'Campanas'),
    (5, 'Autogestion'),
    (6, 'D. Bienvenida'),
    (7, 'D. Stock'),
    (8, 'D. Masiva'),
    (9, 'D. Satisfechos'),
    (10, 'D. Microseguro'),
    (11, 'D. Cancelaciones'),
    (12, 'Asesores'),
    (13, 'Iniciativas'),
    (14, 'Evidencias'),
    (15, 'Capacitaciones'),
    (16, 'Monitoreo'),
    (17, 'Cap. 2'),
    (18, 'Contactab.'),
    (19, 'Telefonia'),
    (20, 'Proyeccion'),
    (21, 'Estrategia'),
    (22, 'Cierre'),
]

# ── Data per slide for editable overlays ──
def get_slide_data(n):
    """Returns (metrics_list, observations_list) for editable overlay"""
    metrics = []
    obs = []
    
    if n == 2:  # Ventas
        total = sum(D.get('ventasLiq', [0]*6))
        cp = sum(D.get('ventasCP', [0]*6))
        vol = sum(D.get('ventasVOL', [0]*6))
        metrics = [
            ('Ventas totales', fmt(total)),
            ('Cuota Protegida', fmt(cp)),
            ('Combo Vida', fmt(vol)),
            ('Meta Vanti', fmt(sum(D.get('metaVanti', [0]*6)))),
            ('Cumpl. Vanti', '147%'),
            ('Meta Xuma', fmt(sum(D.get('metaXuma', [0]*6)))),
            ('Cumpl. Xuma', '96%'),
        ]
        obs = ['El semestre cerro con 13.917 polizas. Vanti al 147 % de su meta. Xuma al 96 %. Liquidacion oficial supera el tablero operativo por ventas subsanadas post-corte.']
    elif n == 3:  # Bases
        metrics = [
            ('Registros', fmt(sum(D.get('registros', [0]*6)))),
            ('Aptos', fmt(sum(D.get('aptos', [0]*6)))),
            ('Ventas Liq.', fmt(sum(D.get('ventasLiq', [0]*6)))),
        ]
        obs = ['Embudo completo: registros -> rechazados -> aptos -> contactados -> ventas. Ver tabla mensual para detalle por mes.']
    elif n == 4:  # Campanas
        metrics = [('Total descartes', fmt(sum(m.get('cant',0) for m in J.get('DESCARTE_MOTIVOS',[]))))]
        obs = ['Distribucion de motivos de rechazo en la base. Cuota Protegida Activa lidera el descarte estructural.']
    elif n == 5:  # Autogestion
        am = J.get('AUTOGESTION_MESES', [])
        if am:
            total_v = sum(m.get('ventas',0) for m in am)
            metrics = [('Ventas autogestion', fmt(total_v))]
        obs = ['Base de financiaciones autogestionadas en plataforma Vanti. Meta ideal: 20 % conversion. Focalizada en 3 asesores desde 26-mar.']
    elif n == 12:  # Asesores
        asesores = J.get('ASESORES', [])
        roster = J.get('ROSTER', {})
        rp = roster.get('rosterPorMes', []) if roster else []
        top5 = sum(sorted([sum(v for v in (a.get('meses',[]) or []) if v is not None) for a in asesores], reverse=True)[:5])
        metrics = [
            ('Total asesores', str(len(asesores))),
            ('Roster actual', str(rp[-1]) if rp else '21'),
            ('Top 5 aporte', fmt(top5)),
        ]
        obs = ['Top 5 aporta el 37 % del semestre (4.279 polizas). Roster crecio de 15 a 21 asesores (+40 %). Incentivos: 8 eventos Vanti + 41 extras Xuma.']
    elif n in (6,7,8,9,10,11):  # Detalles de campana
        det_keys = list(J.get('detalles', {}).keys())
        idx = n - 6
        if 0 <= idx < len(det_keys):
            det = J.get('detalles', {}).get(det_keys[idx], {})
            data = det.get('data', {}) if det else {}
            hl = data.get('highlights', [])
            observaciones = data.get('observaciones', [])
            for h in hl[:3]:
                metrics.append((h.get('label',''), h.get('val','')))
            for o in observaciones:
                obs.append(clean(o))
    elif n == 13:  # Iniciativas
        iniciativas = J.get('INICIATIVAS_1S', [])
        for ini in iniciativas[:3]:
            obs.append(f"{ini.get('nombre','')}: {ini.get('estado','')}")
    elif n == 15:  # Capacitaciones
        caps = J.get('CAPACITACIONES_1S', [])
        metrics = [('Capacitaciones', str(len(caps)))]
        for cap in caps[:3]:
            obs.append(f"{cap.get('mes','')}: {cap.get('tema','')}")
    elif n == 18:  # Contactab
        con = D.get('contactabilidad', [])
        efe = D.get('efectividad', [])
        metrics = [
            ('Contactabilidad prom', f"{sum(con)/len(con):.1f}%" if con else 'S/D'),
            ('Efectividad prom', f"{sum(efe)/len(efe):.1f}%" if efe else 'S/D'),
        ]
    elif n == 20:  # Proyeccion
        total = sum(D.get('ventasLiq', [0]*6))
        metrics = [
            ('Ventas 1S', fmt(total)),
            ('Promedio/mes', fmt(total//6) if total else '0'),
        ]
        obs = ['Meta 2S: 18.000 polizas (3.000/mes) = +29 %. Apalancadores: roster 21 a 30, nuevas campanas, optimizacion de bases.']
    elif n == 21:  # Estrategia
        estr = J.get('ESTRATEGIA_INICIATIVAS', [])
        for e in estr[:3]:
            obs.append(f"{e.get('nombre','')}: {e.get('timeline','')}")
    
    return metrics, obs

# ── SERVER ──
class Handler(SimpleHTTPRequestHandler):
    def log_message(self, *a):
        pass
handler = partial(Handler, directory=DASH_DIR)
server = HTTPServer(('localhost', 0), handler)
port = server.server_address[1]
print(f"Servidor: http://localhost:{port}")
t = threading.Thread(target=server.serve_forever, daemon=True)
t.start()
time.sleep(0.5)

# ── SCREENSHOTS ──
print("Capturando 23 slides...")
screenshots = {}
with sync_playwright() as pw:
    browser = pw.chromium.launch(headless=True)
    page = browser.new_page(viewport={'width': 1280, 'height': 720})
    for num, label in SLIDES:
        try:
            page.goto(f"http://localhost:{port}/televentas/index.html?slide={num}", wait_until='networkidle', timeout=15000)
            time.sleep(0.5)
            # Hide nav chrome
            page.evaluate("""() => {
                ['nav','next-btn','home-btn'].forEach(id => {
                    const el = document.getElementById(id);
                    if (el) el.style.display = 'none';
                });
                const p = document.querySelector('.progress-bar');
                if (p) p.style.display = 'none';
            }""")
            time.sleep(0.2)
            path = os.path.join(OUT_DIR, f'slide_{num:02d}.png')
            page.screenshot(path=path)
            screenshots[num] = path
            print(f"  {num:2d} {label}")
        except Exception as e:
            print(f"  {num:2d} {label} ERROR: {e}")
    browser.close()

# ── BUILD PPT ──
print("\nConstruyendo PPT con capas editables...")
prs = Presentation()
prs.slide_width = PPT_W
prs.slide_height = PPT_H

for num, label in SLIDES:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo = screenshot
    img = screenshots.get(num)
    if img and os.path.exists(img):
        s.shapes.add_picture(img, 0, 0, PPT_W, PPT_H)
    
    # Panel editable inferior con TODOS los datos
    metrics, obs = get_slide_data(num)
    has_content = metrics or obs
    
    if has_content:
        bar_top = Inches(6.75)
        bar_h = Inches(0.70)
        
        # Fondo semitransparente oscuro
        bar = s.shapes.add_shape(1, Inches(0), bar_top, PPT_W, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = DARK
        bar.line.fill.background()
        
        # Construir texto editable
        lines = []
        if metrics:
            parts = [f"{lab}: {val}" for lab, val in metrics]
            lines.append(" | ".join(parts))
        for o in obs:
            lines.append(o)
        txt = "\n".join(lines)
        
        tb = s.shapes.add_textbox(Inches(0.25), bar_top + Inches(0.02), PPT_W - Inches(0.5), bar_h - Inches(0.04))
        tf = tb.text_frame; tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(1) if i > 0 else Pt(0)
            p.space_after = Pt(0)
            r = p.add_run()
            r.text = line
            r.font.size = Pt(6.5) if i == 0 else Pt(6)
            r.font.color.rgb = WHITE
            r.font.name = 'Raleway'
            r.font.bold = (i == 0)
    
    # Titulo editable en barra superior
    tb2 = s.shapes.add_textbox(Inches(0.25), Inches(0.08), Inches(8), Inches(0.35))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(10)
    r2.font.color.rgb = DARK
    r2.font.name = 'Raleway'
    r2.font.bold = True

out = os.path.join(DIR, 'INFORME_TELEVENTAS_HIBRIDO.pptx')
prs.save(out)
print(f"\nPPT generado: {out}")
print(f"{len(prs.slides)} slides con fondo + datos editables en barra inferior")
