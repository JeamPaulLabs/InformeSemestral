#!/usr/bin/env python3
"""PPT esqueleto v3: fondo sin texto + capas editables en posicion exacta"""

import json, os, time, threading, re
from http.server import HTTPServer, SimpleHTTPRequestHandler
from functools import partial
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from playwright.sync_api import sync_playwright

DIR = os.path.dirname(os.path.abspath(__file__))
DASH_DIR = os.path.join(DIR, '..')
OUT_DIR = os.path.join(DIR, 'ppt_slides')
os.makedirs(OUT_DIR, exist_ok=True)

PPT_W = Inches(13.333)
PPT_H = Inches(7.5)
SX = 13.333 / 1280.0
SY = 7.5 / 720.0

def parse_color(s):
    m = re.match(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', str(s))
    return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3))) if m else None

class H(SimpleHTTPRequestHandler):
    def log_message(self, *a): pass
server = HTTPServer(('localhost', 0), partial(H, directory=DASH_DIR))
port = server.server_address[1]
print(f"Server: http://localhost:{port}")
t = threading.Thread(target=server.serve_forever, daemon=True); t.start()
time.sleep(0.5)

prs = Presentation()
prs.slide_width = PPT_W; prs.slide_height = PPT_H

SLIDES = [
    (0, 'Portada', None), (1, 'Cap. 1', None),
    (2, 'Ventas', [('Vanti', "vtasTab('vanti')"), ('Xuma', "vtasTab('xuma')")]),
    (3, 'Bases', None), (4, 'Campanas', None), (5, 'Autogestion', None),
    (6, 'D. Bienvenida', None), (7, 'D. Stock', None), (8, 'D. Masiva', None),
    (9, 'D. Satisfechos', None), (10, 'D. Microseguro', None), (11, 'D. Cancelaciones', None),
    (12, 'Asesores', [('Vanti', "asesoresTab('vanti')"), ('Xuma', "asesoresTab('xuma')")]),
    (13, 'Iniciativas', None), (14, 'Evidencias', None), (15, 'Capacitaciones', None),
    (16, 'Monitoreo', None), (17, 'Cap. 2', None),
    (18, 'Contactab.', [('Mes', "contactabTab('mes')"), ('Campana', "contactabTab('campana')")]),
    (19, 'Telefonia', [('Resumen', "telefoniaTab('resumen')"), ('Zonas', "telefoniaTab('zonas')")]),
    (20, 'Proyeccion', [('Calculo', "proyeccionTab('calc')"), ('Escenario', "proyeccionTab('escenario')")]),
    (21, 'Estrategia', [('Iniciativas', "estrategiaTab('ini')"),
                        ('Cronograma', "estrategiaTab('cron')"),
                        ('KPIs', "estrategiaTab('kpi')")]),
    (22, 'Cierre', None),
]

print("Procesando slides...")

with sync_playwright() as pw:
    browser = pw.chromium.launch(headless=True)
    ctx = browser.new_context(viewport={'width': 1280, 'height': 720})
    page = ctx.new_page()
    
    for num, label, tabs in SLIDES:
        try:
            page.goto(f"http://localhost:{port}/televentas/index.html?slide={num}", wait_until='networkidle', timeout=15000)
            time.sleep(0.8)
            
            # Force canvas to exact position: position:relative so offsetParent chain works
            page.evaluate("""() => {
                const s = document.getElementById('scaler');
                if (s) { s.style.position = 'relative'; s.style.transform = 'scale(1)'; s.style.top = '0'; s.style.left = '0'; }
                document.body.style.margin = '0'; document.body.style.overflow = 'hidden';
            }""")
            time.sleep(0.1)
            
            # Hide floating nav buttons
            page.evaluate("""() => {
                ['prev-btn','next-btn','home-btn','nav'].forEach(id => {
                    const el = document.getElementById(id);
                    if (el) el.style.setProperty('display', 'none', 'important');
                });
                const p = document.querySelector('.progress-bar');
                if (p) p.style.setProperty('display', 'none', 'important');
            }""")
            
            variants = [('', None)]
            if tabs:
                variants = [(f' {s}', js) for s, js in tabs]
            
            for v_suffix, v_js in variants:
                slide_label = label + v_suffix
                
                # Switch tab if needed
                if v_js:
                    page.evaluate(v_js)
                    time.sleep(0.3)
                
                # ── Extract text positions using bounding rects ──
                text_data = page.evaluate("""() => {
                    const scaler = document.getElementById('scaler');
                    if (!scaler) return [];
                    const sRect = scaler.getBoundingClientRect();
                    const scaleX = sRect.width / 1280;
                    const scaleY = sRect.height / 720;
                    const active = document.querySelector('#scaler .slide.active') || scaler;
                    const results = [];
                    const seen = new Set();
                    
                    function canvasRect(el) {
                        const r = el.getBoundingClientRect();
                        return {
                            x: (r.left - sRect.left) / scaleX,
                            y: (r.top - sRect.top) / scaleY,
                            w: r.width / scaleX,
                            h: r.height / scaleY
                        };
                    }
                    
                    // Get all text-containing leaf elements (direct text children)
                    const allEls = active.querySelectorAll('*');
                    allEls.forEach(el => {
                        const style = window.getComputedStyle(el);
                        if (style.display === 'none' || style.visibility === 'hidden') return;
                        if (el.closest('svg') || el.closest('use') || el.tagName === 'use') return;
                        
                        // Direct text only
                        let dt = '';
                        for (const c of el.childNodes) {
                            if (c.nodeType === 3) dt += c.textContent;
                        }
                        const text = dt.trim();
                        if (!text || text.length < 1) return;
                        
                        const cr = canvasRect(el);
                        if (cr.w < 3 || cr.h < 3) return;
                        
                        // Deduplicate by position + text
                        const key = `${cr.x.toFixed(0)},${cr.y.toFixed(0)},${text}`;
                        if (seen.has(key)) return;
                        seen.add(key);
                        
                        const fs = parseFloat(style.fontSize) || 10;
                        const fw = style.fontWeight;
                        
                        results.push({
                            text, x: cr.x, y: cr.y, w: cr.w, h: cr.h,
                            fontSize: fs,
                            bold: fw === 'bold' || parseInt(fw) >= 600,
                            color: style.color,
                            textAlign: style.textAlign || 'left'
                        });
                    });
                    return results;
                }""")
                
                # Deduplicate
                text_data.sort(key=lambda t: -len(t['text']))
                filtered = []
                regions = []
                for t in text_data:
                    rx, ry, rw, rh = t['x'], t['y'], t['w'], t['h']
                    contained = False
                    for ux, uy, uw, uh in regions:
                        if rx >= ux and ry >= uy and rx+rw <= ux+uw and ry+rh <= uy+uh:
                            contained = True; break
                    if not contained:
                        filtered.append(t)
                        regions.append((rx, ry, rw, rh))
                
                # ── Screenshot SKELETON (text hidden, icons preserved) ──
                page.evaluate("""() => {
                    // Save SVG/icon colors
                    window._pptSvgColors = [];
                    document.querySelectorAll('svg, .lucide-ico').forEach(el => {
                        const cs = window.getComputedStyle(el);
                        window._pptSvgColors.push({ el, color: cs.color });
                    });
                    
                    // Hide text on non-SVG elements
                    document.querySelectorAll('*').forEach(el => {
                        if (el.closest('svg') || el.classList.contains('lucide-ico')) return;
                        el.style.setProperty('color', 'transparent', 'important');
                    });
                    
                    // Restore SVG colors
                    window._pptSvgColors.forEach(({ el, color }) => {
                        el.style.setProperty('color', color, 'important');
                    });
                }""")
                time.sleep(0.1)
                
                ss_path = os.path.join(OUT_DIR, f'sk_{num:02d}{v_suffix.replace(" ","_")}.png')
                page.screenshot(path=ss_path)
                
                # Restore text
                page.evaluate("""() => {
                    document.querySelectorAll('*').forEach(el => el.style.removeProperty('color'));
                }""")
                
                # DEBUG: print first 15 coords for Ventas
                if num == 2 and v_suffix == ' Vanti':
                    for t_in in text_data[:20]:
                        print(f"    DEBUG raw: x={t_in['x']:.1f} y={t_in['y']:.1f} w={t_in['w']:.1f} h={t_in['h']:.1f} fs={t_in.get('fontSize',0)} txt={t_in['text'][:40]}")
                    print(f"    DEBUG filtered: {len(filtered)} of {len(text_data)} total")
                
                # ── Build PPT slide ──
                s = prs.slides.add_slide(prs.slide_layouts[6])
                
                # Skeleton background
                if os.path.exists(ss_path):
                    s.shapes.add_picture(ss_path, 0, 0, PPT_W, PPT_H)
                
                # Editable text boxes at exact canvas positions
                for t in filtered:
                    tx, ty, tw, th = t['x'], t['y'], t['w'], t['h']
                    if tw < 5 or th < 3: continue
                    
                    left = tx * SX
                    top = ty * SY
                    width = max(tw * SX, Inches(0.2))
                    height = max(th * SY, Inches(0.1))
                    
                    tb = s.shapes.add_textbox(left, top, width, height)
                    tf = tb.text_frame
                    tf.word_wrap = False
                    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
                    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
                    
                    p = tf.paragraphs[0]
                    p.alignment = {'left': PP_ALIGN.LEFT, 'right': PP_ALIGN.RIGHT,
                                   'center': PP_ALIGN.CENTER}.get(t['textAlign'], PP_ALIGN.LEFT)
                    
                    fs = Pt(min(t['fontSize'] * 0.75, 16))
                    r = p.add_run()
                    r.text = t['text']
                    r.font.size = fs
                    r.font.bold = t['bold']
                    r.font.name = 'Raleway'
                    
                    clr = parse_color(t['color'])
                    if clr: r.font.color.rgb = clr
                    
                    tb.fill.background()
                
                print(f"  {num:2d} {slide_label:25s} -> {len(filtered):3d} textos")
                
        except Exception as e:
            import traceback
            print(f"  {num:2d} {label}: ERROR {e}")
            traceback.print_exc()
    
    browser.close()

out = os.path.join(DIR, 'INFORME_TELEVENTAS_ESQUELETO.pptx')
prs.save(out)
print(f"\nPPT: {out}")
print(f"{len(prs.slides)} slides")
