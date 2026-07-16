#!/usr/bin/env python3
"""PPT hibrido: esqueleto visual (sin texto) + capas editables en posicion exacta"""

import json, os, time, threading, re
from http.server import HTTPServer, SimpleHTTPRequestHandler
from functools import partial
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from playwright.sync_api import sync_playwright

DIR = os.path.dirname(os.path.abspath(__file__))
DASH_DIR = os.path.join(DIR, '..')
OUT_DIR = os.path.join(DIR, 'ppt_slides')
os.makedirs(OUT_DIR, exist_ok=True)

PPT_W = Inches(13.333)
PPT_H = Inches(7.5)
# Canvas 1280x720 -> PPT inches
SX = 13.333 / 1280.0
SY = 7.5 / 720.0

def px(x, y, w=None, h=None):
    l = x * SX
    t = y * SY
    width = w * SX if w else None
    height = h * SY if h else None
    return (l, t, width, height)

SLIDES = [
    (0, 'Portada'), (1, 'Cap. 1'), (2, 'Ventas'), (3, 'Bases'),
    (4, 'Campanas'), (5, 'Autogestion'), (6, 'D. Bienvenida'), (7, 'D. Stock'),
    (8, 'D. Masiva'), (9, 'D. Satisfechos'), (10, 'D. Microseguro'),
    (11, 'D. Cancelaciones'), (12, 'Asesores'), (13, 'Iniciativas'),
    (14, 'Evidencias'), (15, 'Capacitaciones'), (16, 'Monitoreo'),
    (17, 'Cap. 2'), (18, 'Contactab.'), (19, 'Telefonia'),
    (20, 'Proyeccion'), (21, 'Estrategia'), (22, 'Cierre'),
]

# ── SERVER ──
class Handler(SimpleHTTPRequestHandler):
    def log_message(self, *a): pass
handler = partial(Handler, directory=DASH_DIR)
server = HTTPServer(('localhost', 0), handler)
port = server.server_address[1]
print(f"Servidor: http://localhost:{port}")
t = threading.Thread(target=server.serve_forever, daemon=True)
t.start()
time.sleep(0.5)

# ── Helper: parse color from computed style ──
def parse_color(rgb_str):
    """Convert 'rgb(r,g,b)' or 'rgba(r,g,b,a)' to RGBColor or None"""
    m = re.match(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', str(rgb_str))
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    return None

def parse_font_size(s):
    m = re.match(r'([\d.]+)', str(s))
    return float(m.group(1)) if m else 10

# ── SCREENSHOT SKELETONS + EXTRACT TEXT ──
print("Capturando esqueletos y extrayendo textos...")

prs = Presentation()
prs.slide_width = PPT_W
prs.slide_height = PPT_H

with sync_playwright() as pw:
    browser = pw.chromium.launch(headless=True)
    page = browser.new_page(viewport={'width': 1280, 'height': 720})
    
    for num, label in SLIDES:
        print(f"  Slide {num:2d} {label}...", end=' ')
        try:
            page.goto(f"http://localhost:{port}/televentas/index.html?slide={num}", wait_until='networkidle', timeout=15000)
            time.sleep(1.0)  # Let slide animations finish
            
            # Hide nav chrome
            page.evaluate("""() => {
                ['nav','next-btn','home-btn'].forEach(id => {
                    const el = document.getElementById(id);
                    if (el) el.style.display = 'none';
                });
                const p = document.querySelector('.progress-bar');
                if (p) p.style.display = 'none';
            }""")
            time.sleep(0.2)
            
            # Force canvas to exact 1280x720 with no transform offset
            page.evaluate("""() => {
                const scaler = document.getElementById('scaler');
                if (scaler) {
                    scaler.style.transform = 'scale(1)';
                    scaler.style.top = '0';
                    scaler.style.left = '0';
                }
                document.body.style.margin = '0';
                document.body.style.overflow = 'hidden';
            }""")
            time.sleep(0.1)
            
            # ── Extract text positions ──
            text_data = page.evaluate("""() => {
                const results = [];
                const active = document.querySelector('#scaler .slide.active') || document.querySelector('#scaler');
                if (!active) return [];
                
                const allEls = active.querySelectorAll('*');
                allEls.forEach(el => {
                    const style = window.getComputedStyle(el);
                    if (style.display === 'none' || style.visibility === 'hidden') return;
                    if (el.closest('svg')) return;
                    
                    let directText = '';
                    for (const child of el.childNodes) {
                        if (child.nodeType === 3) directText += child.textContent;
                    }
                    const text = directText.trim();
                    if (!text || text.length < 2) return;
                    
                    const r = el.getBoundingClientRect();
                    if (r.width < 3 || r.height < 3) return;
                    
                    const fs = parseFloat(style.fontSize) || 10;
                    const fw = style.fontWeight;
                    
                    results.push({
                        text, x: r.left, y: r.top, w: r.width, h: r.height,
                        fontSize: fs,
                        bold: fw === 'bold' || parseInt(fw) >= 600,
                        color: style.color,
                        textAlign: style.textAlign || 'left'
                    });
                });
                
                if (results.length === 0) {
                    const fullText = active.innerText || active.textContent || '';
                    if (fullText.trim()) {
                        results.push({
                            text: fullText.trim().substring(0, 500),
                            x: 40, y: 80, w: 1200, h: 500,
                            fontSize: 12, bold: false,
                            color: 'rgb(0,0,0)', textAlign: 'left'
                        });
                    }
                }
                return results;
            }""")
            
            # Deduplicate overlapping regions: keep the one with more text
            text_data.sort(key=lambda t: -len(t['text']))
            filtered = []
            regions = []
            for t in text_data:
                rx, ry, rw, rh = t['x'], t['y'], t['w'], t['h']
                contained = False
                for ux, uy, uw, uh in regions:
                    if rx >= ux and ry >= uy and rx+rw <= ux+uw and ry+rh <= uy+uh:
                        contained = True; break
                if not contained:
                    filtered.append(t)
                    regions.append((rx, ry, rw, rh))
            
            # ── Screenshot skeleton (text hidden via CSS, preserves layout & icons) ──
            page.evaluate("""() => {
                // Save original colors of SVG/icon elements to restore later
                window.__pptIconColors = [];
                document.querySelectorAll('svg, svg *, .lucide-ico, [class*="ico"]').forEach(el => {
                    const cs = window.getComputedStyle(el);
                    if (cs.color && cs.color !== 'transparent') {
                        window.__pptIconColors.push({ el, color: el.style.color || '' });
                    }
                });
                
                // Hide all text by making color transparent
                // But first save original body color for restoration
                window.__pptBodyColor = document.body.style.color || '';
                document.querySelectorAll('*').forEach(el => {
                    el.style.setProperty('color', 'transparent', 'important');
                });
                
                // Restore SVG/icon element colors
                window.__pptIconColors.forEach(({ el, color }) => {
                    el.style.setProperty('color', color || '#120180', 'important');
                });
                
                // Also restore images (they might use color for alt text)
                document.querySelectorAll('img').forEach(el => {
                    el.style.removeProperty('color');
                });
            }""")
            time.sleep(0.1)
            
            skeleton_path = os.path.join(OUT_DIR, f'skeleton_{num:02d}.png')
            page.screenshot(path=skeleton_path)
            
            # Restore text visibility
            page.evaluate("""() => {
                document.querySelectorAll('*').forEach(el => {
                    el.style.removeProperty('color');
                });
                if (window.__pptBodyColor) {
                    document.body.style.color = window.__pptBodyColor;
                }
            }""")
            
            # ── Build PPT slide ──
            s = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Skeleton as full background
            if os.path.exists(skeleton_path):
                s.shapes.add_picture(skeleton_path, 0, 0, PPT_W, PPT_H)
            
            # Editable text boxes at exact positions
            for t in filtered:
                tx = t['x']
                ty = t['y']
                tw = t['w']
                th = t['h']
                
                # Convert to inches
                left = tx * SX
                top = ty * SY
                width = max(tw * SX, Inches(0.3))
                height = max(th * SY, Inches(0.15))
                
                # Skip very small elements
                if tw < 5 or th < 3:
                    continue
                
                tb = s.shapes.add_textbox(left, top, width, height)
                tf = tb.text_frame
                tf.word_wrap = True
                tf.margin_left = Pt(0)
                tf.margin_right = Pt(0)
                tf.margin_top = Pt(0)
                tf.margin_bottom = Pt(0)
                
                p = tf.paragraphs[0]
                p.alignment = {
                    'left': PP_ALIGN.LEFT, 'right': PP_ALIGN.RIGHT,
                    'center': PP_ALIGN.CENTER, 'justify': PP_ALIGN.JUSTIFY
                }.get(t['textAlign'], PP_ALIGN.LEFT)
                
                # Parse font size
                fs = Pt(min(t['fontSize'] * 0.75, 14))  # scale for PPT
                
                r = p.add_run()
                r.text = t['text']
                r.font.size = fs
                r.font.bold = t['bold']
                r.font.name = 'Raleway'
                
                clr = parse_color(t['color'])
                if clr:
                    r.font.color.rgb = clr
                
                # Make background transparent so skeleton shows
                tb.fill.background()
            
            print(f"OK ({len(filtered)} textos)")
            
        except Exception as e:
            print(f"ERROR: {e}")
    
    browser.close()

# ── Save ──
out = os.path.join(DIR, 'INFORME_TELEVENTAS_ESQUELETO.pptx')
prs.save(out)
print(f"\nPPT generado: {out}")
print(f"{len(prs.slides)} slides con esqueleto + textos editables en posicion exacta")
